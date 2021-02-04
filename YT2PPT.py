from pytube import YouTube
from pptx import Presentation
from pptx.util import Inches
import os
import shutil
import cv2
import skimage.measure as measure

url = "https://www.youtube.com/watch?v=bNb2fEVKeEo"
epsilon = 0.015

# Create Local Dir
dirName = 'Frames'
try:
    # Create target Directory
    os.mkdir(dirName)
    print("Directory ", dirName, " Created ")
except FileExistsError:
    print("Directory ", dirName, " already exists")

# Downloading Video to Local Dir
pdfImg = []
ytvideo = YouTube(url)
frameRate = 1  # //it will capture image in each 0.5 second
video = ytvideo.streams.first().download(filename="videotoextract")
print("VIDEO Successfully download")
print("Start cutting into Frames")


def compare_img(img1, img2):
    """
    Compare two images with scikit-image measure.
    :param img1: first image
    :param img2: second image
    :return: Score according to the difference between both images
    """
    # load the two input images
    image_a = cv2.imread("Frames\\" + img1)
    image_b = cv2.imread("Frames\\" + img2)
    # convert the images to grayscale
    gray_a = cv2.cvtColor(image_a, cv2.COLOR_BGR2GRAY)
    gray_b = cv2.cvtColor(image_b, cv2.COLOR_BGR2GRAY)
    # compute the Structural Similarity Index (SSIM) between the two
    # images, ensuring that the difference image is returned
    (score, diff) = measure.compare_ssim(gray_a, gray_b, full=True)
    return score


def delete_image(index):
    """
    Delete a picture from Frames path
    :param index: index of the image (int)
    """
    filename = 'image' + index + ".jpg"
    file_path = os.path.join("Frames", filename)
    try:
        if os.path.isfile(file_path) or os.path.islink(file_path):
            os.remove(file_path)
        elif os.path.isdir(file_path):
            shutil.rmtree(file_path)
    except Exception as e:
        print('Failed to delete %s. Reason: %s' % (file_path, e))


vid_cap = cv2.VideoCapture('videotoextract.mp4')


def video_to_frames(seconds):
    """
    Cut the current video into frames and export frames in Frames dir
    :param seconds: Second in the video to start cutting into frames
    :return: boolean
    """
    global count
    vid_cap.set(cv2.CAP_PROP_POS_MSEC, seconds * 1000)
    has_frames, image = vid_cap.read()
    if has_frames:
        cv2.imwrite("Frames\\image" + str(count) + ".jpg", image)  # save frame as JPG file
        if count > 1:
            if compare_img("image" + str(count - 1) + ".jpg", "image" + str(count) + ".jpg") > 1 - epsilon:
                delete_image(str(count))
                count = count - 1
                if count not in pdfImg:
                    if len(pdfImg) > 0:
                        bool = 0
                        for i in range(len(pdfImg)):
                            if (compare_img("image" + str(pdfImg[i]) + ".jpg",
                                            "image" + str(count) + ".jpg") > 1 - epsilon):
                                bool = 1
                        if bool == 0:
                            pdfImg.append(count)
                    elif len(pdfImg) == 0:
                        pdfImg.append(count)

    return has_frames


count = 1
sec = 0

# Cut video into frames and keep only ppt ones
success = video_to_frames(sec)
while success:
    count = count + 1
    sec = sec + frameRate
    sec = round(sec, 2)
    success = video_to_frames(sec)

vid_cap.release()
cv2.destroyAllWindows()
print("Successfully cut into Frames")
print("Creating ppt")

# Creating the ppt:
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
blank_slide_layout = prs.slide_layouts[6]

left = top = 0
for i in range(0, len(pdfImg)):
    slide = prs.slides.add_slide(blank_slide_layout)
    slide.shapes.add_picture("Frames\\image" + str(pdfImg[i]) + ".jpg", left, top, height=prs.slide_height)

prs.save('YouTubeSlides.pptx')
print("Ppt was saved")
print("Cleaning Env")
# Clean Directory:
shutil.rmtree('Frames')
os.remove('videotoextract.mp4')
